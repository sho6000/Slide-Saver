import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import datetime
import streamlit as st

# App title and description
st.set_page_config(page_title="Slide Saver", layout="wide")
st.title("🎤 Slide Saver - Generate Lyrics Presentation")
st.markdown("Customize and generate a PowerPoint presentation from your song lyrics.")

def generate_presentation(lyrics):
    with st.expander("🎨 Font Customization", expanded=True):
        col1, col2 = st.columns(2)
        font_name = col1.selectbox("Font:", ["Arial", "Times New Roman", "Special Elite", "Helvetica", "AnonymousPro", "Agrandir"], help="Default is Arial")
        font_size = col2.slider("Font Size:", min_value=10, max_value=72, value=32)
        
        col1, col2 = st.columns(2)
        font_color = col1.color_picker("Font Color:", "#FFFFFF")
        font_color_rgb = tuple(int(font_color[i:i+2], 16) for i in (1, 3, 5))

        bold = col2.checkbox("Bold")
        italics = col1.checkbox("Italics")
        underline = col2.checkbox("Underline")

    with st.expander("📏 Line Spacing", expanded=False):
        line_spacing_options = [1.0, 1.5, 2.0, 2.5, 3.0]
        line_spacing = st.radio("Select Line Spacing:", line_spacing_options, horizontal=True)

    with st.expander("🖼️ Background", expanded=False):
        background_choice = st.radio("Select Background:", ["Black", "White", "Custom Image"], horizontal=True)
        custom_background = None
        if background_choice == "Custom Image":
            custom_background = st.file_uploader("Upload Custom Background", type=["jpg", "jpeg", "png"])
    
    with st.expander("📐 Text Alignment", expanded=True):
        alignment_options = {"Left": PP_ALIGN.LEFT, "Center": PP_ALIGN.CENTER, "Right": PP_ALIGN.RIGHT}
        text_alignment = st.radio("Choose Text Alignment:", list(alignment_options.keys()), horizontal=True)

    # Generate button with styling
    st.markdown("---")
    st.markdown('<div style="text-align:center">', unsafe_allow_html=True)
    if st.button("🎬 Generate Presentation", help="Click to create your PowerPoint slides"):
        presentation = Presentation()
        stanzas = lyrics.split("\n\n")

        for stanza_text in stanzas:
            if not stanza_text.strip():
                continue

            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            if background_choice == "Black":
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
            elif background_choice == "White":
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            elif background_choice == "Custom Image" and custom_background:
                slide.shapes.add_picture(custom_background, Inches(0), Inches(0), Inches(10), Inches(7.5))

            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            for line in stanza_text.split('\n'):
                if not line.strip():
                    continue
                p = text_frame.add_paragraph()
                p.text = line
                run = p.runs[0]
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*font_color_rgb)
                p.alignment = alignment_options[text_alignment]  # Apply selected alignment
                if bold: run.font.bold = True
                if italics: run.font.italic = True
                if underline: run.font.underline = True
                text_frame.space_after = Pt(font_size * line_spacing - font_size)

        file_name = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        presentation.save(file_name)

        with open(file_name, "rb") as file:
            st.download_button("📥 Download Presentation", file.read(), file_name, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

        os.remove(file_name)
        st.success("✅ Presentation generated successfully!")

    st.markdown("</div>", unsafe_allow_html=True)
